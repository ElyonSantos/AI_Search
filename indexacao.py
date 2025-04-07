import zipfile  # Para manipular arquivos ZIP
import rarfile  # Para manipular arquivos RAR
import tarfile  # Para manipular arquivos TAR
import py7zr    # Para manipular arquivos 7Z
import csv # Faz suporte para arquivos do Excel
import json # Para extensões json
import os  # Biblioteca para manipulação de diretórios e arquivos
import sqlite3  # Banco de dados SQLite para armazenar informações dos arquivos
import datetime  # Para registrar data e hora da modificação dos arquivos
import subprocess # Faz parte da transcrição de audio/video
from PyPDF2 import PdfReader  # Biblioteca para leitura de arquivos PDF
import pytesseract  # OCR para extrair texto de imagens
from PIL import Image  # Manipulação de imagens
from docx import Document  # Para ler arquivos DOCX
import openpyxl  # Para ler arquivos XLSX
import tensorflow as tf # Importa modelo pré-treinado para identificar as coisas nas imagens
import speech_recognition as sr # Importa as coisas necessárias para transcrever audio
from pydub import AudioSegment # Tambem possui as coisas para transcrever audio
import xml.etree.ElementTree as ET # Para extensões xml
from pptx import Presentation # Suporte para arquivos pptx
from sentence_transformers import SentenceTransformer # Faz os embeddings para o projeto
import numpy as np # Para carregar os embeddings


# Tirar avisos de bibliotecas
os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"
tf.compat.v1.logging.set_verbosity(tf.compat.v1.logging.ERROR)

# Carrega o modelo para gerar embeddings
modelo_embedding = SentenceTransformer('all-mpnet-base-v2')

print(f"Modelo carregado: {modelo_embedding}")
print("Dimensão do embedding gerado:", modelo_embedding.get_sentence_embedding_dimension())

# Define os caminhos para o FFmpeg e FFprobe
AudioSegment.converter = r"C:\ffmpeg\bin\ffmpeg.exe"
AudioSegment.ffprobe = r"C:\ffmpeg\bin\ffprobe.exe"

# Configura o caminho para o executável do Tesseract
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Define o caminho da pasta e do banco de dados
DATABASE_DIR = "database"
DATABASE_PATH = os.path.join(DATABASE_DIR, "index_arquivos.db")


def conectar_banco(nome_db=DATABASE_PATH):
    if not os.path.exists(DATABASE_DIR):
        os.makedirs(DATABASE_DIR)
        print(f"Pasta '{DATABASE_DIR}' criada!")

    try:
        # Define um timeout de 10 segundos para evitar bloqueios prolongados
        conn = sqlite3.connect(nome_db, timeout=10)
        cursor = conn.cursor()

        # Otimizações para evitar bloqueios
        cursor.execute("PRAGMA journal_mode=WAL;")  # Permite múltiplos acessos simultâneos
        cursor.execute("PRAGMA synchronous = NORMAL;")  # Melhora desempenho de escrita
        cursor.execute("PRAGMA temp_store = MEMORY;")  # Usa RAM para tabelas temporárias
        cursor.execute("PRAGMA foreign_keys = ON;")  # Garante integridade referencial

        # Criar tabelas se não existirem
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS arquivos (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                nome TEXT,
                caminho TEXT UNIQUE,
                tipo TEXT,
                tamanho INTEGER,
                data_modificacao TEXT,
                status TEXT DEFAULT 'ativo'
            )
        ''')

        cursor.execute('''
            CREATE TABLE IF NOT EXISTS conteudo (
                id INTEGER PRIMARY KEY,
                texto TEXT,
                FOREIGN KEY (id) REFERENCES arquivos(id) ON DELETE CASCADE
            )
        ''')

        cursor.execute('''
            CREATE TABLE IF NOT EXISTS embeddings (
                id INTEGER PRIMARY KEY,
                embedding_nome BLOB,
                embedding_caminho BLOB,
                embedding_tipo BLOB,
                embedding_tamanho BLOB,
                embedding_data BLOB,
                embedding_status BLOB,
                embedding_conteudo BLOB,
                FOREIGN KEY (id) REFERENCES arquivos(id) ON DELETE CASCADE
            )
        ''')

        conn.commit()
        return conn, cursor

    except sqlite3.Error as e:
        print(f"Erro ao conectar ao banco de dados: {e}")
        return None, None  # Retorna valores vazios se houver erro


def gerar_embedding(valor):
    if not valor:
        return None  # Retorna None se o valor estiver vazio
    try:
        # Converte o valor para string e gera o embedding
        embedding = modelo_embedding.encode(str(valor)).flatten()
        if np.isnan(embedding).any():
            return None  # Retorna None se houver NaN
        return embedding
    except Exception as e:
        print(f"Erro ao gerar embedding: {e}")
        return None


def salvar_embedding(cursor, id_arquivo, coluna_embedding, embedding):
    if embedding is None:
        print(f"Embedding inválido para {coluna_embedding} do arquivo {id_arquivo}")
        return
    try:
        # Serializa o embedding
        embedding_serializado = np.array(embedding, dtype=np.float32).tobytes()

        # Garante que o registro existe na tabela 'embeddings'
        cursor.execute("INSERT OR IGNORE INTO embeddings (id) VALUES (?)", (id_arquivo,))

        # Atualiza a coluna específica
        cursor.execute(f'''
            UPDATE embeddings 
            SET {coluna_embedding} = ? 
            WHERE id = ?
        ''', (embedding_serializado, id_arquivo))
    except Exception as e:
        print(f"Erro ao salvar embedding para ID={id_arquivo} ({coluna_embedding}): {e}")


def atualizar_banco_embeddings_nome(cursor):
    cursor.execute("SELECT id, nome FROM arquivos")
    for id_arquivo, nome in cursor.fetchall():
        embedding = gerar_embedding(nome)
        salvar_embedding(cursor, id_arquivo, "embedding_nome", embedding)


def atualizar_banco_embeddings_caminho(cursor):
    cursor.execute("SELECT id, caminho FROM arquivos")
    for id_arquivo, caminho in cursor.fetchall():
        embedding = gerar_embedding(caminho)
        salvar_embedding(cursor, id_arquivo, "embedding_caminho", embedding)


def atualizar_banco_embeddings_tipo(cursor):
    cursor.execute("SELECT id, tipo FROM arquivos")
    for id_arquivo, tipo in cursor.fetchall():
        embedding = gerar_embedding(tipo)
        salvar_embedding(cursor, id_arquivo, "embedding_tipo", embedding)


def atualizar_banco_embeddings_tamanho(cursor):
    cursor.execute("SELECT id, tamanho FROM arquivos")
    for id_arquivo, tamanho in cursor.fetchall():
        embedding = gerar_embedding(tamanho)
        salvar_embedding(cursor, id_arquivo, "embedding_tamanho", embedding)


def atualizar_banco_embeddings_data(cursor):
    cursor.execute("SELECT id, data_modificacao FROM arquivos")
    for id_arquivo, data in cursor.fetchall():
        embedding = gerar_embedding(data)
        salvar_embedding(cursor, id_arquivo, "embedding_data", embedding)


def atualizar_banco_embeddings_status(cursor):
    cursor.execute("SELECT id, status FROM arquivos")
    for id_arquivo, status in cursor.fetchall():
        embedding = gerar_embedding(status)
        salvar_embedding(cursor, id_arquivo, "embedding_status", embedding)


def atualizar_banco_embeddings_conteudo(cursor):
    cursor.execute("SELECT id, texto FROM conteudo")
    for id_arquivo, texto in cursor.fetchall():
        embedding = gerar_embedding(texto)
        salvar_embedding(cursor, id_arquivo, "embedding_conteudo", embedding)


def atualizar_todos_embeddings():
    conn, cursor = conectar_banco()
    print("Atualizando embeddings...")

    try:
        cursor.execute('''
            INSERT OR IGNORE INTO embeddings (id)
            SELECT id FROM arquivos
            WHERE id NOT IN (SELECT id FROM embeddings)
        ''')

        atualizar_banco_embeddings_nome(cursor)
        atualizar_banco_embeddings_caminho(cursor)
        atualizar_banco_embeddings_tipo(cursor)
        atualizar_banco_embeddings_tamanho(cursor)
        atualizar_banco_embeddings_data(cursor)
        atualizar_banco_embeddings_status(cursor)
        atualizar_banco_embeddings_conteudo(cursor)

        conn.commit()
    except sqlite3.OperationalError as e:
        print(f"Erro no banco de dados: {e}")
    finally:
        conn.close()  # Fecha a conexão corretamente


def extrair_audio_ffmpeg(caminho_arquivo):
    try:
        print(f"Convertendo arquivo para MP3 usando FFmpeg: {caminho_arquivo}")
        caminho_mp3 = f"{caminho_arquivo.rsplit('.', 1)[0]}.mp3"
        comando = [
            "ffmpeg",
            "-loglevel", "error",  # Reduz o nível de verbosidade
            "-i", caminho_arquivo,
            "-q:a", "0",
            "-map", "a",
            caminho_mp3
        ]
        subprocess.run(comando, check=True)

        # Verifica se o arquivo MP3 foi criado
        if os.path.exists(caminho_mp3):
            print(f"Áudio extraído para: {caminho_mp3}")
            return caminho_mp3
        else:
            print(f"Falha ao criar arquivo MP3: {caminho_mp3}")
            return None
    except subprocess.CalledProcessError as e:
        print(f"Erro ao converter arquivo {caminho_arquivo} com FFmpeg: {e}")
        return None


def transcrever_audio(caminho_arquivo):
    try:
        print(f"Processando áudio: {caminho_arquivo}")

        # Converte o arquivo para WAV (se necessário)
        if not caminho_arquivo.lower().endswith(".wav"):
            print("Convertendo arquivo para WAV...")
            audio = AudioSegment.from_file(caminho_arquivo)
            caminho_wav = caminho_arquivo.rsplit(".", 1)[0] + ".wav"
            audio.export(caminho_wav, format="wav")
            caminho_arquivo = caminho_wav
            print(f"Arquivo convertido para: {caminho_wav}")

            # Verifica se o arquivo WAV foi criado
            if not os.path.exists(caminho_wav):
                print(f"Falha ao criar arquivo WAV: {caminho_wav}")
                return "Não foi possível transcrever o áudio."

        # Usa o SpeechRecognition para transcrever o áudio
        print("Transcrevendo áudio...")
        recognizer = sr.Recognizer()
        with sr.AudioFile(caminho_arquivo) as source:
            audio_data = recognizer.record(source)
            texto = recognizer.recognize_google(audio_data, language="pt-BR")
            print(f"Texto transcrito: {texto}")
            return texto

    except FileNotFoundError as e:
        print(f"Arquivo não encontrado: {caminho_arquivo}. Erro: {e}")
        return "Não foi possível transcrever o áudio."
    except sr.UnknownValueError:
        print(f"Não foi possível entender o áudio: {caminho_arquivo}")
        return "Não foi possível transcrever o áudio."
    except Exception as e:
        print(f"Erro ao processar áudio {caminho_arquivo}: {e}")
        return "Não foi possível transcrever o áudio."


# Função para identificar objetos em uma imagem
# noinspection PyUnresolvedReferences
def identificar_objeto(caminho_imagem):
    # Carrega um modelo pré-treinado (MobileNetV2)
    modelo = tf.keras.applications.MobileNetV2(weights='imagenet')

    # Pré-processa a imagem
    imagem = tf.keras.preprocessing.image.load_img(caminho_imagem, target_size=(224, 224))  # Redimensiona para 224x224
    imagem_array = tf.keras.preprocessing.image.img_to_array(imagem)
    imagem_array = tf.keras.applications.mobilenet_v2.preprocess_input(imagem_array)
    imagem_array = tf.expand_dims(imagem_array, axis=0)  # Adiciona dimensão batch

    # Faz a previsão
    previsoes = modelo.predict(imagem_array)
    resultados = tf.keras.applications.mobilenet_v2.decode_predictions(previsoes, top=5)[0]

    # Retorna os resultados como texto descritivo
    descricao = ", ".join([f"{nome_objeto} ({probabilidade:.2f})" for _, nome_objeto, probabilidade in resultados])
    return descricao


# Função para extrair texto de um PDF e armazená-lo no banco
def extrair_conteudo_pdf(caminho, cursor, id):
    try:
        with open(caminho, "rb") as f:
            reader = PdfReader(f)  # Abre o PDF
            texto_completo = "\n".join([page.extract_text() for page in reader.pages if
                                        page.extract_text()])  # Extrai texto de todas as páginas
        for i in range(0, len(texto_completo), 5000):  # Divide em partes de 5000 caracteres
            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)",
                           (id, texto_completo[i:i + 5000]))  # Insere no banco
    except Exception as e:
        print(f"Erro ao processar PDF {caminho}: {e}")


# Função para extrair .txt
def extrair_texto_txt(caminho):
    try:
        with open(caminho, "r", encoding="utf-8", errors="ignore") as arquivo:
            texto_completo = arquivo.read()
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar TXT {caminho}: {e}")
        return ""


# Função para extrair texto de imagens usando OCR
def extrair_texto_imagem(caminho):
    try:
        return pytesseract.image_to_string(Image.open(caminho))  # Retorna texto extraído da imagem
    except Exception as e:
        print(f"Erro ao processar imagem {caminho}: {e}")
        return ""  # Retorna string vazia se houver erro


# Função para extrair texto de arquivos DOCX
def extrair_texto_docx(caminho):
    try:
        doc = Document(caminho)  # Usa a classe Document do módulo docx
        texto_completo = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar DOCX {caminho}: {e}")
        return ""


# Função para extrair texto de arquivos XLSX
def extrair_texto_xlsx(caminho):
    try:
        workbook = openpyxl.load_workbook(caminho)
        texto_completo = ""
        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows(values_only=True):
                texto_completo += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar XLSX {caminho}: {e}")
        return ""


def extrair_texto_csv(caminho):
    try:
        with open(caminho, "r", encoding="utf-8") as arquivo:
            leitor = csv.reader(arquivo)
            texto_completo = "\n".join([", ".join(linha) for linha in leitor])
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar CSV {caminho}: {e}")
        return ""


def extrair_texto_json(caminho):
    try:
        with open(caminho, "r", encoding="utf-8") as arquivo:
            dados = json.load(arquivo)
            texto_completo = json.dumps(dados, ensure_ascii=False, indent=4)
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar JSON {caminho}: {e}")
        return ""


def extrair_texto_xml(caminho):
    try:
        tree = ET.parse(caminho)
        root = tree.getroot()
        texto_completo = "".join(root.itertext())
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar XML {caminho}: {e}")
        return ""


def extrair_texto_pptx(caminho):
    try:
        presentation = Presentation(caminho)
        texto_completo = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto_completo += shape.text + "\n"
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar PPTX {caminho}: {e}")
        return ""


def extrair_texto_bak(caminho):
    try:
        with open(caminho, "rb") as arquivo:
            conteudo = arquivo.read()
            texto_completo = conteudo.decode("utf-8", errors="ignore")
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar BAK {caminho}: {e}")
        return ""


def extrair_texto_log(caminho):
    try:
        with open(caminho, "r", encoding="utf-8", errors="ignore") as arquivo:
            texto_completo = arquivo.read()
        return texto_completo
    except Exception as e:
        print(f"Erro ao processar LOG {caminho}: {e}")
        return ""


def extrair_e_indexar_zip(caminho, cursor, id):
    try:
        with zipfile.ZipFile(caminho, 'r') as zip_ref:
            temp_dir = os.path.join(os.path.dirname(caminho), f"temp_{id}")
            os.makedirs(temp_dir, exist_ok=True)
            zip_ref.extractall(temp_dir)

            for root, _, files in os.walk(temp_dir):
                for nome_arquivo in files:
                    caminho_completo = os.path.join(root, nome_arquivo)
                    extensao = os.path.splitext(nome_arquivo)[-1].lower()

                    if deve_ignorar_arquivo(nome_arquivo):
                        continue

                    tamanho = os.path.getsize(caminho_completo)
                    data_mod = datetime.datetime.fromtimestamp(
                        os.path.getmtime(caminho_completo)).isoformat()

                    cursor.execute('''
                        INSERT OR REPLACE INTO arquivos (nome, caminho, tipo, tamanho, data_modificacao, status)
                        VALUES (?, ?, ?, ?, ?, 'ativo')
                    ''', (nome_arquivo, caminho_completo, extensao, tamanho, data_mod))
                    novo_id = cursor.lastrowid

                    if extensao == ".pdf":
                        extrair_conteudo_pdf(caminho_completo, cursor, novo_id)
                    elif extensao in [".jpg", ".png", ".gif"]:
                        texto_extraido = extrair_texto_imagem(caminho_completo)
                        descricao_objetos = identificar_objeto(caminho_completo)
                        texto_final = f"Texto extraído: {texto_extraido}\nObjetos identificados: {descricao_objetos}"
                        cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto_final))
                    elif extensao == ".docx":
                        texto = extrair_texto_docx(caminho_completo)
                        if texto:
                            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto))
                    elif extensao == ".xlsx":
                        texto = extrair_texto_xlsx(caminho_completo)
                        if texto:
                            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto))
    except Exception as e:
        print(f"Erro ao processar ZIP {caminho}: {e}")


def extrair_e_indexar_tar(caminho, cursor, id):
    try:
        with tarfile.open(caminho, 'r:*') as tar_ref:
            temp_dir = os.path.join(os.path.dirname(caminho), f"temp_{id}")
            os.makedirs(temp_dir, exist_ok=True)
            tar_ref.extractall(temp_dir)

            for root, _, files in os.walk(temp_dir):
                for nome_arquivo in files:
                    caminho_completo = os.path.join(root, nome_arquivo)
                    extensao = os.path.splitext(nome_arquivo)[-1].lower()

                    if deve_ignorar_arquivo(nome_arquivo):
                        continue

                    tamanho = os.path.getsize(caminho_completo)
                    data_mod = datetime.datetime.fromtimestamp(
                        os.path.getmtime(caminho_completo)).isoformat()

                    cursor.execute('''
                        INSERT OR REPLACE INTO arquivos (nome, caminho, tipo, tamanho, data_modificacao, status)
                        VALUES (?, ?, ?, ?, ?, 'ativo')
                    ''', (nome_arquivo, caminho_completo, extensao, tamanho, data_mod))
                    novo_id = cursor.lastrowid

                    if extensao == ".pdf":
                        extrair_conteudo_pdf(caminho_completo, cursor, novo_id)
                    elif extensao in [".jpg", ".png", ".gif"]:
                        texto_extraido = extrair_texto_imagem(caminho_completo)
                        descricao_objetos = identificar_objeto(caminho_completo)
                        texto_final = f"Texto extraído: {texto_extraido}\nObjetos identificados: {descricao_objetos}"
                        cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto_final))
                    elif extensao == ".docx":
                        texto = extrair_texto_docx(caminho_completo)
                        if texto:
                            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto))
                    elif extensao == ".xlsx":
                        texto = extrair_texto_xlsx(caminho_completo)
                        if texto:
                            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto))
    except Exception as e:
        print(f"Erro ao processar TAR {caminho}: {e}")


def extrair_e_indexar_rar(caminho, cursor, id):
    try:
        with rarfile.RarFile(caminho, 'r') as rar_ref:
            temp_dir = os.path.join(os.path.dirname(caminho), f"temp_{id}")
            os.makedirs(temp_dir, exist_ok=True)
            rar_ref.extractall(temp_dir)

            for root, _, files in os.walk(temp_dir):
                for nome_arquivo in files:
                    caminho_completo = os.path.join(root, nome_arquivo)
                    extensao = os.path.splitext(nome_arquivo)[-1].lower()

                    if deve_ignorar_arquivo(nome_arquivo):
                        continue

                    tamanho = os.path.getsize(caminho_completo)
                    data_mod = datetime.datetime.fromtimestamp(
                        os.path.getmtime(caminho_completo)).isoformat()

                    cursor.execute('''
                        INSERT OR REPLACE INTO arquivos (nome, caminho, tipo, tamanho, data_modificacao, status)
                        VALUES (?, ?, ?, ?, ?, 'ativo')
                    ''', (nome_arquivo, caminho_completo, extensao, tamanho, data_mod))
                    novo_id = cursor.lastrowid

                    if extensao == ".pdf":
                        extrair_conteudo_pdf(caminho_completo, cursor, novo_id)
                    elif extensao in [".jpg", ".png", ".gif"]:
                        texto_extraido = extrair_texto_imagem(caminho_completo)
                        descricao_objetos = identificar_objeto(caminho_completo)
                        texto_final = f"Texto extraído: {texto_extraido}\nObjetos identificados: {descricao_objetos}"
                        cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto_final))
                    elif extensao == ".docx":
                        texto = extrair_texto_docx(caminho_completo)
                        if texto:
                            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto))
                    elif extensao == ".xlsx":
                        texto = extrair_texto_xlsx(caminho_completo)
                        if texto:
                            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto))
    except Exception as e:
        print(f"Erro ao processar RAR {caminho}: {e}")


def extrair_e_indexar_7z(caminho, cursor, id):
    try:
        with py7zr.SevenZipFile(caminho, mode='r') as seven_ref:
            temp_dir = os.path.join(os.path.dirname(caminho), f"temp_{id}")
            os.makedirs(temp_dir, exist_ok=True)
            seven_ref.extractall(path=temp_dir)

            for root, _, files in os.walk(temp_dir):
                for nome_arquivo in files:
                    caminho_completo = os.path.join(root, nome_arquivo)
                    extensao = os.path.splitext(nome_arquivo)[-1].lower()

                    if deve_ignorar_arquivo(nome_arquivo):
                        continue

                    tamanho = os.path.getsize(caminho_completo)
                    data_mod = datetime.datetime.fromtimestamp(
                        os.path.getmtime(caminho_completo)).isoformat()

                    cursor.execute('''
                        INSERT OR REPLACE INTO arquivos (nome, caminho, tipo, tamanho, data_modificacao, status)
                        VALUES (?, ?, ?, ?, ?, 'ativo')
                    ''', (nome_arquivo, caminho_completo, extensao, tamanho, data_mod))
                    novo_id = cursor.lastrowid

                    if extensao == ".pdf":
                        extrair_conteudo_pdf(caminho_completo, cursor, novo_id)
                    elif extensao in [".jpg", ".png", ".gif"]:
                        texto_extraido = extrair_texto_imagem(caminho_completo)
                        descricao_objetos = identificar_objeto(caminho_completo)
                        texto_final = f"Texto extraído: {texto_extraido}\nObjetos identificados: {descricao_objetos}"
                        cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto_final))
                    elif extensao == ".docx":
                        texto = extrair_texto_docx(caminho_completo)
                        if texto:
                            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto))
                    elif extensao == ".xlsx":
                        texto = extrair_texto_xlsx(caminho_completo)
                        if texto:
                            cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (novo_id, texto))
    except Exception as e:
        print(f"Erro ao processar 7Z {caminho}: {e}")


# Função para verificar se um arquivo deve ser ignorado
def deve_ignorar_arquivo(nome_arquivo):
    """
    Verifica se um arquivo deve ser ignorado com base no nome ou extensão.

    Args:
        nome_arquivo (str): Nome completo do arquivo.

    Returns:
        bool: True se o arquivo deve ser ignorado, False caso contrário.
    """
    # Lista de padrões de nomes de arquivos a serem ignorados
    padroes_ignorados = [
        "~$",  # Arquivos temporários do Word
        ".tmp",  # Arquivos temporários genéricos
        ".DS_Store",  # Arquivos de sistema macOS
        "Thumbs.db"  # Arquivos de miniaturas do Windows
    ]

    # Lista de extensões indesejadas
    extensoes_ignoradas = [
        ".tmp", ".DS_Store", ".db", ".ini", ".sys"
    ]

    # Verifica se o nome do arquivo contém algum dos padrões ignorados
    for padrao in padroes_ignorados:
        if padrao in nome_arquivo:
            return True

    # Verifica se a extensão do arquivo está na lista de extensões ignoradas
    _, extensao = os.path.splitext(nome_arquivo)
    if extensao.lower() in extensoes_ignoradas:
        return True

    # Se nenhum critério for atendido, o arquivo não deve ser ignorado
    return False


# Função principal para percorrer diretórios e indexar arquivos
def indexar_diretorio(diretorio):
    conn, cursor = conectar_banco()  # Conecta ao banco
    total_arquivos = sum([len(files) for _, _, files in os.walk(diretorio)])  # Conta o total de arquivos
    print(f"Total de arquivos no diretório: {total_arquivos}")

    contador = 0  # Inicializa o contador de arquivos processados
    novos_arquivos = 0  # Contador de novos arquivos indexados
    arquivos_existentes_no_banco = set()  # Conjunto para rastrear arquivos no banco
    arquivos_atualmente_no_diretorio = set()  # Conjunto para rastrear arquivos no diretório
    arquivos_ignorados = 0  # Contador de arquivos ignorados

    # Obtém todos os arquivos já registrados no banco
    cursor.execute("SELECT caminho FROM arquivos")
    for linha in cursor.fetchall():
        arquivos_existentes_no_banco.add(linha[0])

    for raiz, _, arquivos in os.walk(diretorio):  # Percorre todos os arquivos do diretório
        for arquivo in arquivos:
            # Verifica se o arquivo deve ser ignorado
            if deve_ignorar_arquivo(arquivo):
                arquivos_ignorados += 1  # Incrementa o contador de arquivos ignorados
                continue  # Ignora o arquivo

            contador += 1  # Incrementa o contador
            caminho_completo = os.path.join(raiz, arquivo)  # Obtém caminho completo

            # Garante que o caminho seja sempre uma string
            caminho_completo = str(caminho_completo)

            tamanho = os.path.getsize(caminho_completo)  # Obtém tamanho do arquivo
            data_mod = datetime.datetime.fromtimestamp(
                os.path.getmtime(caminho_completo)).isoformat()  # Obtém data de modificação
            extensao = os.path.splitext(arquivo)[-1].lower()  # Obtém a extensão do arquivo

            arquivos_atualmente_no_diretorio.add(caminho_completo)  # Adiciona ao conjunto

            # Verifica se o arquivo já está no banco de dados
            cursor.execute("SELECT data_modificacao FROM arquivos WHERE caminho = ?", (caminho_completo,))
            resultado = cursor.fetchone()

            if resultado:
                data_mod_banco = resultado[0]
                if data_mod_banco == data_mod:  # Se a data de modificação é igual, ignora o arquivo
                    continue

            # Se o arquivo não está no banco ou foi modificado, insere/atualiza
            cursor.execute('''
                INSERT OR REPLACE INTO arquivos (nome, caminho, tipo, tamanho, data_modificacao, status)
                VALUES (?, ?, ?, ?, ?, 'ativo')
            ''', (arquivo, caminho_completo, extensao, tamanho, data_mod))  # Insere metadados do arquivo no banco
            id = cursor.lastrowid  # Obtém o ID do arquivo inserido

            # Processa o conteúdo do arquivo, dependendo do tipo
            if extensao == ".pdf":
                extrair_conteudo_pdf(caminho_completo, cursor, id)  # Processa PDF

            elif extensao in [".jpg", ".png", ".gif"]:
                # Extrai texto da imagem usando Tesseract
                texto_extraido = extrair_texto_imagem(caminho_completo)

                # Identifica objetos na imagem usando visão computacional
                try:
                    descricao_objetos = identificar_objeto(caminho_completo)
                except Exception as e:
                    print(f"Erro ao identificar objetos na imagem {caminho_completo}: {e}")
                    descricao_objetos = "Não foi possível identificar objetos."

                # Combina o texto extraído e a descrição dos objetos
                texto_final = f"Texto extraído: {texto_extraido}\nObjetos identificados: {descricao_objetos}"

                # Insere o texto combinado no banco de dados
                cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)",
                               (id, texto_final))
            elif extensao == ".docx":
                texto = extrair_texto_docx(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)",
                                   (id, texto))
            elif extensao == ".xlsx":
                texto = extrair_texto_xlsx(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)",
                                   (id, texto))
            elif extensao in [".mp3", ".wav", ".flac"]:
                # Processa áudio
                descricao = transcrever_audio(caminho_completo)
                cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)",
                               (id, descricao))
            elif extensao in [".mp4", ".avi", ".mkv", ".mov"]:
                # Processa vídeo
                print(f"Processando vídeo: {caminho_completo}")
                caminho_audio = extrair_audio_ffmpeg(caminho_completo)
                if caminho_audio:
                    descricao = transcrever_audio(caminho_audio)
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)",
                                   (id, descricao))
                else:
                    print(f"Falha ao processar vídeo: {caminho_completo}")

            elif extensao == ".csv":
                texto = extrair_texto_csv(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (id, texto))

            elif extensao == ".json":
                texto = extrair_texto_json(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (id, texto))

            elif extensao == ".xml":
                texto = extrair_texto_xml(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (id, texto))

            elif extensao == ".pptx":
                texto = extrair_texto_pptx(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (id, texto))

            elif extensao == ".bak":
                texto = extrair_texto_bak(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (id, texto))

            elif extensao == ".log":
                texto = extrair_texto_log(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (id, texto))
            elif extensao == ".txt":
                texto = extrair_texto_txt(caminho_completo)
                if texto:
                    cursor.execute("INSERT INTO conteudo (id, texto) VALUES (?, ?)", (id, texto))

            #Sessão abaixo dos diferentões:
            #a indexação deles é grande, por isso é feito separado

            elif extensao == ".zip":
                extrair_e_indexar_zip(caminho_completo, cursor, id)
            elif extensao == ".rar":
                extrair_e_indexar_rar(caminho_completo, cursor, id)
            elif extensao == ".tar":
                extrair_e_indexar_tar(caminho_completo, cursor, id)
            elif extensao == ".7z":
                extrair_e_indexar_7z(caminho_completo, cursor, id)

            novos_arquivos += 1  # Incrementa o contador de novos arquivos

            # Atualiza o progresso
            mensagem = f"\033[33mIndexando arquivos ({contador}/{total_arquivos})\n\033[0m"
            print("\r" + " " * len(mensagem), end="")  # Limpa a linha anterior
            print("\r" + mensagem, end="", flush=True)  # Imprime a nova mensagem

    # Marca arquivos ausentes como "excluídos"
    arquivos_excluidos = arquivos_existentes_no_banco - arquivos_atualmente_no_diretorio
    for caminho_excluido in arquivos_excluidos:
        cursor.execute("UPDATE arquivos SET status = 'excluido' WHERE caminho = ?", (caminho_excluido,))

    conn.commit()  # Salva mudanças no banco
    conn.close()  # Fecha a conexão

    # Atualiza todos os embeddings no banco de dados
    print("\nAtualizando banco de dados com todos os embeddings...")
    atualizar_todos_embeddings()

    # Exibe resumo final
    print(f"\n\033[34mIndexação concluída! {novos_arquivos} novos arquivos foram indexados.\033[0m")
    if arquivos_ignorados > 0:
        print(f"\033[34m{arquivos_ignorados} arquivos foram ignorados.\033[0m")
    if len(arquivos_excluidos) > 0:
        print(f"\033[34m{len(arquivos_excluidos)} arquivos foram marcados como excluídos.\033[0m")
    else:
        print("\n\033[34mNenhum novo arquivo foi encontrado para indexação.\033[0m")


def main():
    diretorio = input("Informe o diretório para indexação: ")  # Solicita diretório ao usuário
    if os.path.isdir(diretorio):  # Verifica se o diretório existe
        indexar_diretorio(diretorio)  # Chama a função de indexação

    else:
        print("Diretório inválido!")


if __name__ == "__main__":
    main()  # Inicia o programa